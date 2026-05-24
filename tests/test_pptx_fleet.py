"""Test fleet for enhanced PPTX generation with PlantUML diagrams and speaker notes.

Runs a battery of prompts against the document generation pipeline to validate:
  1. PlantUML rendering (Java + JAR check)
  2. Enhanced markdown parsing (diagrams, notes)
  3. Full PPTX build with embedded images
  4. Fallback behavior when PlantUML is unavailable

Usage:
    python -m tests.test_pptx_fleet
    # or
    pytest tests/test_pptx_fleet.py -v
"""

from __future__ import annotations

import io
import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "backend"))

from app.plantuml import find_plantuml_jar, is_available, render_puml_to_png
from app.pptx_builder import build_pptx, parse_enhanced_markdown


# ---------------------------------------------------------------------------
# Infrastructure checks
# ---------------------------------------------------------------------------


class TestPlantUMLInfrastructure:
    """Verify Java and PlantUML JAR are available."""

    def test_java_installed(self):
        import subprocess

        result = subprocess.run(["java", "-version"], capture_output=True)
        assert result.returncode == 0, "Java is not installed"

    def test_plantuml_jar_exists(self):
        jar = find_plantuml_jar()
        assert jar is not None, f"PlantUML JAR not found at ~/.local/lib/plantuml.jar"
        assert jar.stat().st_size > 1_000_000, "JAR file seems too small"

    def test_is_available(self):
        assert is_available(), "PlantUML is not available (need Java + JAR)"

    def test_render_simple_diagram(self):
        puml = "@startuml\nAlice -> Bob: Hello\n@enduml"
        png = render_puml_to_png(puml)
        assert png is not None, "Failed to render simple diagram"
        assert len(png) > 100, "PNG output too small"
        assert png[:8] == b"\x89PNG\r\n\x1a\n", "Output is not a valid PNG"


# ---------------------------------------------------------------------------
# Markdown parsing tests
# ---------------------------------------------------------------------------


class TestEnhancedMarkdownParsing:
    """Test the enhanced markdown parser handles diagrams and notes."""

    def test_basic_slides(self):
        md = "# Title\n- subtitle\n## Slide 2\n- bullet 1\n- bullet 2"
        slides = parse_enhanced_markdown(md)
        assert len(slides) == 2
        assert slides[0]["title"] == "Title"
        assert slides[1]["bullets"] == ["bullet 1", "bullet 2"]

    def test_diagram_block(self):
        md = (
            "# My Deck\n- overview\n"
            "## Architecture\n"
            "```plantuml\n@startuml\nA -> B\n@enduml\n```\n"
        )
        slides = parse_enhanced_markdown(md)
        assert len(slides) == 2
        assert "@startuml" in slides[1]["diagram"]
        assert "A -> B" in slides[1]["diagram"]

    def test_speaker_notes(self):
        md = (
            "# Title\n- sub\n"
            "<!-- notes: Welcome everyone to this presentation -->\n"
            "## Content\n- point\n"
            "<!-- notes: This slide covers the main points -->\n"
        )
        slides = parse_enhanced_markdown(md)
        assert "Welcome everyone" in slides[0]["notes"]
        assert "main points" in slides[1]["notes"]

    def test_mixed_content(self):
        md = (
            "# House Document System\n- A comprehensive overview\n"
            "<!-- notes: Introduce the system -->\n"
            "## How Upload Works\n"
            "```plantuml\n@startuml\nstart\n:Upload file;\n:Extract text;\nstop\n@enduml\n```\n"
            "<!-- notes: Walk through the upload flow -->\n"
            "## Key Features\n- Search\n- Ask AI\n- Generate docs\n"
            "<!-- notes: Highlight the three main features -->\n"
        )
        slides = parse_enhanced_markdown(md)
        assert len(slides) == 3
        assert slides[0]["notes"] == "Introduce the system"
        assert "@startuml" in slides[1]["diagram"]
        assert slides[1]["notes"] == "Walk through the upload flow"
        assert len(slides[2]["bullets"]) == 3


# ---------------------------------------------------------------------------
# PPTX build tests
# ---------------------------------------------------------------------------


class TestPPTXBuild:
    """Test PPTX generation produces valid files."""

    def test_basic_build(self):
        slides = [
            {"title": "Test Deck", "bullets": ["A test presentation"], "notes": "Welcome", "diagram": ""},
            {"title": "Content", "bullets": ["Point 1", "Point 2"], "notes": "Discuss points", "diagram": ""},
        ]
        pptx_bytes = build_pptx(slides)
        assert len(pptx_bytes) > 1000
        # Verify it's a valid ZIP (PPTX is ZIP-based)
        assert pptx_bytes[:2] == b"PK"

    def test_build_with_diagram(self):
        slides = [
            {"title": "Diagram Test", "bullets": ["With PlantUML"], "notes": "", "diagram": ""},
            {
                "title": "Workflow",
                "bullets": [],
                "notes": "This shows the workflow",
                "diagram": "@startuml\nstart\n:Step 1;\n:Step 2;\nstop\n@enduml",
            },
        ]
        pptx_bytes = build_pptx(slides)
        assert len(pptx_bytes) > 5000  # Should be larger with embedded image

        # Verify the PPTX contains an image
        from pptx import Presentation

        prs = Presentation(io.BytesIO(pptx_bytes))
        assert len(prs.slides) == 2
        # Check slide 2 has a picture shape
        slide2 = prs.slides[1]
        has_picture = any(shape.shape_type == 13 for shape in slide2.shapes)
        assert has_picture, "Diagram slide should contain an embedded picture"

    def test_build_with_speaker_notes(self):
        slides = [
            {"title": "Notes Test", "bullets": ["sub"], "notes": "Welcome to the deck", "diagram": ""},
            {"title": "Slide 2", "bullets": ["a", "b"], "notes": "Talk about a and b", "diagram": ""},
        ]
        pptx_bytes = build_pptx(slides)

        from pptx import Presentation

        prs = Presentation(io.BytesIO(pptx_bytes))
        notes1 = prs.slides[0].notes_slide.notes_text_frame.text
        notes2 = prs.slides[1].notes_slide.notes_text_frame.text
        assert "Welcome" in notes1
        assert "Talk about" in notes2

    def test_full_pipeline_from_markdown(self):
        """End-to-end: enhanced markdown → parse → build → valid PPTX."""
        md = (
            "# House Document Search Overview\n"
            "- Your personal document assistant\n"
            "<!-- notes: Welcome everyone. Today we'll walk through how the system works. -->\n"
            "\n"
            "## Document Ingestion Pipeline\n"
            "```plantuml\n"
            "@startuml\n"
            "!theme plain\n"
            "start\n"
            ":Upload PDF;\n"
            "if (Has text layer?) then (yes)\n"
            "  :Extract with pypdf;\n"
            "else (no)\n"
            "  :OCR with Bedrock Vision;\n"
            "endif\n"
            ":Classify document;\n"
            ":Index in OpenSearch;\n"
            "stop\n"
            "@enduml\n"
            "```\n"
            "<!-- notes: This diagram shows how documents flow through the system from upload to searchable index. -->\n"
            "\n"
            "## Search & Ask AI\n"
            "- Hybrid search combines keyword and semantic matching\n"
            "- Ask AI uses RAG to answer questions with citations\n"
            "- Results include relevance scores and source snippets\n"
            "<!-- notes: The search system is the core value proposition. -->\n"
            "\n"
            "## Document Generation\n"
            "- Generate new documents grounded in your uploaded content\n"
            "- Supports DOCX, PDF, PPTX, PNG, and plain text\n"
            "- Template filling preserves original formatting\n"
            "<!-- notes: This is the newest feature and what we're demonstrating today. -->\n"
        )
        slides = parse_enhanced_markdown(md)
        assert len(slides) == 4

        pptx_bytes = build_pptx(slides)
        assert pptx_bytes[:2] == b"PK"

        from pptx import Presentation

        prs = Presentation(io.BytesIO(pptx_bytes))
        assert len(prs.slides) == 4

        # Title slide has notes
        assert "Welcome" in prs.slides[0].notes_slide.notes_text_frame.text

        # Diagram slide has a picture
        has_picture = any(shape.shape_type == 13 for shape in prs.slides[1].shapes)
        assert has_picture, "Ingestion pipeline slide should have embedded diagram"


# ---------------------------------------------------------------------------
# Test fleet: prompts that would be sent to Bedrock for PPTX generation
# These validate the format instructions produce parseable output
# ---------------------------------------------------------------------------


class TestPromptFleet:
    """Simulated Bedrock outputs — tests that realistic model responses parse correctly."""

    FLEET = [
        {
            "name": "HOA_overview",
            "description": "Overview of HOA rules and governance",
            "markdown": (
                "# HOA Rules & Governance\n"
                "- Understanding your community guidelines\n"
                "<!-- notes: This presentation summarizes the key HOA rules from your documents. -->\n"
                "\n"
                "## Architectural Review Process\n"
                "```plantuml\n"
                "@startuml\n"
                "!theme plain\n"
                "start\n"
                ":Submit modification request;\n"
                ":ARC reviews within 30 days;\n"
                "if (Approved?) then (yes)\n"
                "  :Begin work within 60 days;\n"
                "else (no)\n"
                "  :Revise and resubmit;\n"
                "endif\n"
                "stop\n"
                "@enduml\n"
                "```\n"
                "<!-- notes: The ARC process typically takes 30 days. All exterior changes need approval. -->\n"
                "\n"
                "## Common Restrictions\n"
                "- Fence height maximum: 6 feet\n"
                "- No commercial vehicles parked overnight\n"
                "- Lawn maintenance required every 2 weeks\n"
                "- Exterior paint colors must be from approved palette\n"
                "<!-- notes: These are the most commonly cited violations in the community. -->\n"
                "\n"
                "## Dues & Assessments\n"
                "- Monthly dues: varies by community\n"
                "- Special assessments require 2/3 vote\n"
                "- Late fees apply after 30 days\n"
                "- Liens can be placed for unpaid dues\n"
                "<!-- notes: Financial obligations are outlined in the CC&Rs. -->\n"
            ),
        },
        {
            "name": "inspection_summary",
            "description": "Home inspection findings presentation",
            "markdown": (
                "# Home Inspection Summary\n"
                "- Key findings and recommendations\n"
                "<!-- notes: This deck summarizes the major findings from your home inspection report. -->\n"
                "\n"
                "## Inspection Workflow\n"
                "```plantuml\n"
                "@startuml\n"
                "!theme plain\n"
                "skinparam backgroundColor #FEFEFE\n"
                "actor Inspector\n"
                "participant House\n"
                "participant Report\n"
                "Inspector -> House: Exterior inspection\n"
                "Inspector -> House: Interior inspection\n"
                "Inspector -> House: Systems check\n"
                "Inspector -> Report: Document findings\n"
                "Report -> Inspector: Generate report\n"
                "@enduml\n"
                "```\n"
                "<!-- notes: The inspector follows a systematic process covering all major systems. -->\n"
                "\n"
                "## Major Findings\n"
                "- Roof: 5-7 years remaining life\n"
                "- HVAC: functioning but aging (15+ years)\n"
                "- Foundation: minor settling, no structural concern\n"
                "- Electrical: panel up to code\n"
                "<!-- notes: No deal-breakers found, but budget for roof replacement in 5 years. -->\n"
                "\n"
                "## Recommended Repairs\n"
                "- Seal gaps around windows (weatherproofing)\n"
                "- Clean gutters and extend downspouts\n"
                "- Service HVAC before winter\n"
                "- Repair grading near foundation\n"
                "<!-- notes: These are low-cost items that prevent bigger problems later. -->\n"
            ),
        },
        {
            "name": "closing_process",
            "description": "Home closing process walkthrough",
            "markdown": (
                "# Home Closing Process\n"
                "- From offer to keys\n"
                "<!-- notes: Let's walk through what happens between your offer being accepted and getting the keys. -->\n"
                "\n"
                "## Closing Timeline\n"
                "```plantuml\n"
                "@startuml\n"
                "!theme plain\n"
                "start\n"
                ":Offer accepted;\n"
                ":Earnest money deposited;\n"
                ":Home inspection (7-10 days);\n"
                ":Appraisal ordered;\n"
                ":Title search;\n"
                ":Final walkthrough;\n"
                ":Closing day - sign documents;\n"
                ":Keys handed over;\n"
                "stop\n"
                "@enduml\n"
                "```\n"
                "<!-- notes: The typical timeline is 30-45 days from accepted offer to closing. -->\n"
                "\n"
                "## Key Documents\n"
                "- Closing Disclosure (review 3 days before)\n"
                "- Deed of Trust / Mortgage\n"
                "- Title Insurance Policy\n"
                "- Homeowner's Insurance Binder\n"
                "- Property Survey\n"
                "<!-- notes: The Closing Disclosure is the most important - compare it to your Loan Estimate. -->\n"
                "\n"
                "## Costs at Closing\n"
                "- Down payment balance\n"
                "- Closing costs (2-5% of purchase price)\n"
                "- Prepaid taxes and insurance\n"
                "- HOA transfer fees\n"
                "<!-- notes: Bring a cashier's check or wire funds the day before. Personal checks are not accepted. -->\n"
            ),
        },
        {
            "name": "insurance_overview",
            "description": "Homeowner's insurance policy overview",
            "markdown": (
                "# Homeowner's Insurance Overview\n"
                "- Understanding your coverage\n"
                "<!-- notes: This presentation breaks down what your homeowner's insurance covers and what it doesn't. -->\n"
                "\n"
                "## Coverage Types\n"
                "- Dwelling coverage (structure)\n"
                "- Personal property (belongings)\n"
                "- Liability protection\n"
                "- Additional living expenses\n"
                "- Medical payments to others\n"
                "<!-- notes: Most policies are HO-3 which covers the structure for all perils except exclusions. -->\n"
                "\n"
                "## Claims Process\n"
                "```plantuml\n"
                "@startuml\n"
                "!theme plain\n"
                "start\n"
                ":Damage occurs;\n"
                ":Document with photos;\n"
                ":File claim with insurer;\n"
                ":Adjuster visits;\n"
                "if (Exceeds deductible?) then (yes)\n"
                "  :Receive payout;\n"
                "else (no)\n"
                "  :Pay out of pocket;\n"
                "endif\n"
                "stop\n"
                "@enduml\n"
                "```\n"
                "<!-- notes: Always document damage immediately with photos before making any repairs. -->\n"
                "\n"
                "## What's NOT Covered\n"
                "- Flood damage (separate policy needed)\n"
                "- Earthquake damage\n"
                "- Normal wear and tear\n"
                "- Pest damage (termites, etc.)\n"
                "<!-- notes: Flood insurance is separate and required if you're in a FEMA flood zone. -->\n"
            ),
        },
        {
            "name": "system_architecture",
            "description": "Document search system architecture (meta - about this app)",
            "markdown": (
                "# Document Search System\n"
                "- Architecture and data flow\n"
                "<!-- notes: This is a technical overview of how the document search application works. -->\n"
                "\n"
                "## System Components\n"
                "```plantuml\n"
                "@startuml\n"
                "!theme plain\n"
                "skinparam backgroundColor #FEFEFE\n"
                "skinparam componentStyle rectangle\n"
                "actor User\n"
                "package \"Frontend\" {\n"
                "  [Svelte App]\n"
                "}\n"
                "package \"Backend\" {\n"
                "  [FastAPI]\n"
                "  [Generator]\n"
                "}\n"
                "database \"PostgreSQL\" as pg\n"
                "database \"OpenSearch\" as os\n"
                "cloud \"AWS Bedrock\" as bedrock\n"
                "User --> [Svelte App]\n"
                "[Svelte App] --> [FastAPI]\n"
                "[FastAPI] --> pg\n"
                "[FastAPI] --> os\n"
                "[Generator] --> bedrock\n"
                "@enduml\n"
                "```\n"
                "<!-- notes: The system uses a standard RAG architecture with OpenSearch for hybrid retrieval. -->\n"
                "\n"
                "## Data Flow\n"
                "- Upload: PDF → extract text → classify → index\n"
                "- Search: query → hybrid search → ranked results\n"
                "- Ask: query → retrieve context → Bedrock → answer with citations\n"
                "- Generate: prompt → retrieve context → Bedrock → format conversion\n"
                "<!-- notes: Every operation is grounded in the user's actual documents. -->\n"
                "\n"
                "## Key Design Decisions\n"
                "- Markdown-first generation (single Bedrock call)\n"
                "- Local format conversion (no additional API costs)\n"
                "- PlantUML diagrams for visual workflows\n"
                "- Speaker notes for presentation context\n"
                "<!-- notes: The markdown-first approach keeps costs low while supporting multiple output formats. -->\n"
            ),
        },
    ]

    @pytest.mark.parametrize("case", FLEET, ids=[c["name"] for c in FLEET])
    def test_parse_fleet_prompt(self, case):
        """Each fleet prompt parses into valid slide data."""
        slides = parse_enhanced_markdown(case["markdown"])
        assert len(slides) >= 3, f"{case['name']}: expected at least 3 slides"
        assert slides[0]["title"], f"{case['name']}: missing title"
        # At least one slide should have notes
        has_notes = any(s.get("notes") for s in slides)
        assert has_notes, f"{case['name']}: no speaker notes found"

    @pytest.mark.parametrize("case", FLEET, ids=[c["name"] for c in FLEET])
    def test_build_fleet_pptx(self, case):
        """Each fleet prompt builds a valid PPTX."""
        slides = parse_enhanced_markdown(case["markdown"])
        pptx_bytes = build_pptx(slides)
        assert pptx_bytes[:2] == b"PK", f"{case['name']}: not a valid PPTX/ZIP"
        assert len(pptx_bytes) > 5000, f"{case['name']}: PPTX too small"

    @pytest.mark.parametrize("case", FLEET, ids=[c["name"] for c in FLEET])
    def test_fleet_diagrams_render(self, case):
        """Diagrams in fleet prompts render to PNG successfully."""
        slides = parse_enhanced_markdown(case["markdown"])
        diagram_slides = [s for s in slides if s.get("diagram")]
        for ds in diagram_slides:
            png = render_puml_to_png(ds["diagram"])
            assert png is not None, f"{case['name']}: diagram failed to render"
            assert png[:8] == b"\x89PNG\r\n\x1a\n", f"{case['name']}: not valid PNG"

    @pytest.mark.parametrize("case", FLEET, ids=[c["name"] for c in FLEET])
    def test_fleet_notes_preserved(self, case):
        """Speaker notes survive the full build pipeline."""
        slides = parse_enhanced_markdown(case["markdown"])
        pptx_bytes = build_pptx(slides)

        from pptx import Presentation

        prs = Presentation(io.BytesIO(pptx_bytes))
        for i, sd in enumerate(slides):
            if sd.get("notes"):
                actual_notes = prs.slides[i].notes_slide.notes_text_frame.text
                assert sd["notes"] in actual_notes, (
                    f"{case['name']} slide {i}: notes not preserved"
                )

    @pytest.mark.parametrize("case", FLEET, ids=[c["name"] for c in FLEET])
    def test_fleet_diagrams_embedded(self, case):
        """Diagram slides contain embedded pictures in the final PPTX."""
        slides = parse_enhanced_markdown(case["markdown"])
        pptx_bytes = build_pptx(slides)

        from pptx import Presentation

        prs = Presentation(io.BytesIO(pptx_bytes))
        for i, sd in enumerate(slides):
            if sd.get("diagram"):
                has_picture = any(shape.shape_type == 13 for shape in prs.slides[i].shapes)
                assert has_picture, (
                    f"{case['name']} slide {i}: diagram not embedded as picture"
                )


# ---------------------------------------------------------------------------
# Integration test: convert_to_pptx from generator module
# ---------------------------------------------------------------------------


class TestGeneratorIntegration:
    """Test that the generator.convert_to_pptx uses the enhanced builder."""

    def test_convert_to_pptx_with_diagram(self):
        from app.generator import convert_to_pptx

        md = (
            "# Test\n- sub\n"
            "## Diagram\n"
            "```plantuml\n@startuml\nA -> B: test\n@enduml\n```\n"
            "## Bullets\n- one\n- two\n"
        )
        result = convert_to_pptx(md)
        assert result[:2] == b"PK"

        from pptx import Presentation

        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 3

    def test_convert_to_pptx_backward_compatible(self):
        """Old-style markdown (no diagrams/notes) still works."""
        from app.generator import convert_to_pptx

        md = "# Old Style Deck\n- subtitle\n## Slide 2\n- bullet a\n- bullet b\n## Slide 3\n- x\n"
        result = convert_to_pptx(md)
        assert result[:2] == b"PK"

        from pptx import Presentation

        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 3


if __name__ == "__main__":
    pytest.main([__file__, "-v", "--tb=short"])
