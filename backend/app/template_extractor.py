"""Template extraction from uploaded documents.

Analyzes a document's structure (not content) and returns a JSON
representation that can be used to generate new documents with the
same layout but different content.
"""

from __future__ import annotations

import json
import logging
import os
from io import BytesIO
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_template(file_bytes: bytes, filename: str) -> dict:
    """Extract structure from a document file.

    Returns a dict describing the document's layout:
    {
        "type": "form" | "report" | "presentation" | "document",
        "title": str,
        "sections": [...],
        "source_format": str,
    }
    """
    ext = Path(filename).suffix.lower()

    # Extract raw text for Bedrock analysis
    text = _get_text_for_analysis(file_bytes, ext)

    # Try Bedrock-enhanced extraction first
    model_id = os.getenv("BEDROCK_TEMPLATE_MODEL_ID", os.getenv("BEDROCK_MODEL_ID", ""))
    if model_id and text and ext in (".docx", ".doc"):
        # For DOCX: always use local extraction (XML is authoritative for structure)
        # Bedrock only adds title/type detection
        local = _extract_docx_template(file_bytes)
        try:
            structure = _bedrock_extract_structure(text, ext, model_id)
            if structure:
                # Take Bedrock's title and type if better
                if structure.get("title") and structure["title"] != "Untitled":
                    local["title"] = structure["title"]
                if structure.get("formatting"):
                    local["formatting"] = structure["formatting"]
        except Exception as e:
            logger.warning("Bedrock template extraction failed: %s", e)
        _fix_bibliography_pattern(local)
        return local
    elif model_id and text:
        try:
            structure = _bedrock_extract_structure(text, ext, model_id)
            if structure and structure.get("sections"):
                if ext == ".pdf":
                    _enrich_pdf_structure(structure, text)
                _fix_bibliography_pattern(structure)
                return structure
        except Exception as e:
            logger.warning("Bedrock template extraction failed, falling back to local: %s", e)

    # Fallback to local extraction
    if ext == ".pdf":
        return _extract_pdf_template(file_bytes)
    elif ext in (".docx", ".doc"):
        return _extract_docx_template(file_bytes)
    elif ext == ".pptx":
        return _extract_pptx_template(file_bytes)
    elif ext == ".md":
        return _extract_md_template(file_bytes.decode("utf-8", errors="ignore"))
    elif ext in (".jpg", ".jpeg", ".png", ".tiff", ".tif"):
        return _extract_image_template(file_bytes, ext)
    else:
        raise ValueError(f"Unsupported template format: {ext}")


def _get_text_for_analysis(file_bytes: bytes, ext: str) -> str:
    """Extract text from file for Bedrock analysis, including font/style metadata."""
    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(BytesIO(file_bytes))
            parts = []
            # Get embedded fonts
            page0 = reader.pages[0] if reader.pages else None
            if page0 and '/Resources' in page0 and '/Font' in page0['/Resources']:
                fonts = page0['/Resources']['/Font']
                font_names = []
                for _, font_obj in fonts.items():
                    resolved = font_obj.get_object() if hasattr(font_obj, 'get_object') else font_obj
                    base_font = str(resolved.get('/BaseFont', '')).lstrip('/')
                    if '+' in base_font:
                        base_font = base_font.split('+', 1)[1]
                    if base_font:
                        font_names.append(base_font)
                if font_names:
                    parts.append(f"[FONTS USED: {', '.join(font_names)}]")
            parts.append(f"[PAGES: {len(reader.pages)}]")
            for pi, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(f"--- Page {pi+1} ---\n{text}")
            return "\n".join(parts)
        elif ext in (".docx", ".doc"):
            return _get_docx_text_with_formatting(file_bytes)
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(BytesIO(file_bytes))
            parts = []
            for si, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_texts.append(para.text.strip())
                if slide_texts:
                    parts.append(f"--- Slide {si+1} ---\n" + "\n".join(slide_texts))
            return "\n".join(parts)
        elif ext == ".md":
            return file_bytes.decode("utf-8", errors="ignore")
    except Exception:
        pass
    return ""


def _get_docx_text_with_formatting(file_bytes: bytes) -> str:
    """Extract DOCX text including SDTs, fonts, styles, and page layout."""
    import zipfile
    from lxml import etree

    parts = []
    ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
          'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

    with zipfile.ZipFile(BytesIO(file_bytes)) as z:
        # Get theme fonts
        theme_files = [n for n in z.namelist() if 'theme' in n and n.endswith('.xml')]
        for tf in theme_files:
            try:
                theme_root = etree.fromstring(z.read(tf))
                major = theme_root.find('.//a:majorFont/a:latin', ns)
                minor = theme_root.find('.//a:minorFont/a:latin', ns)
                if major is not None or minor is not None:
                    parts.append(f"[THEME FONTS: headings={major.get('typeface') if major is not None else 'default'}, body={minor.get('typeface') if minor is not None else 'default'}]")
            except Exception:
                pass

        # Get document defaults from styles
        style_files = [n for n in z.namelist() if 'styles' in n and n.endswith('.xml') and 'glossary' not in n]
        for sf in style_files:
            try:
                styles_root = etree.fromstring(z.read(sf))
                doc_defaults = styles_root.find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}docDefaults')
                if doc_defaults is not None:
                    sz = doc_defaults.find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}sz')
                    if sz is not None:
                        size_pt = int(sz.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val')) / 2
                        parts.append(f"[DEFAULT SIZE: {size_pt}pt]")
            except Exception:
                pass

        # Get page layout from document.xml
        doc_xml = z.read('word/document.xml')
        root = etree.fromstring(doc_xml)

        sect_pr = root.find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}sectPr')
        if sect_pr is not None:
            pg_sz = sect_pr.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pgSz')
            pg_mar = sect_pr.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pgMar')
            if pg_sz is not None:
                w = int(pg_sz.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}w', 0)) / 1440
                h = int(pg_sz.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}h', 0)) / 1440
                parts.append(f"[PAGE SIZE: {w:.1f}x{h:.1f}in]")
            if pg_mar is not None:
                wns = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
                margins = {k: int(pg_mar.get(f'{wns}{k}', 0)) / 1440 for k in ['top', 'bottom', 'left', 'right']}
                parts.append(f"[MARGINS: top={margins['top']:.2f} bottom={margins['bottom']:.2f} left={margins['left']:.2f} right={margins['right']:.2f}in]")

        # Extract content from SDTs (structured document tags) AND paragraphs
        body = root.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}body')
        if body is None:
            return "\n".join(parts)

        wns = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'

        # Walk body children in order
        for child in body:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag

            if tag == 'sdt':
                # Structured document tag
                texts = child.findall(f'.//{wns}t')
                content = ''.join(t.text or '' for t in texts).strip()
                if not content:
                    continue
                # Get style
                pstyle = child.find(f'.//{wns}pStyle')
                style = pstyle.get(f'{wns}val') if pstyle is not None else ""
                # Get run formatting
                rpr = child.find(f'.//{wns}rPr')
                fmt = _parse_run_props(rpr, wns) if rpr is not None else ""
                parts.append(f"[{style}]{fmt} {content}")

            elif tag == 'p':
                texts = child.findall(f'.//{wns}t')
                content = ''.join(t.text or '' for t in texts).strip()
                if not content:
                    continue
                pstyle = child.find(f'.//{wns}pStyle')
                style = pstyle.get(f'{wns}val') if pstyle is not None else ""
                rpr = child.find(f'.//{wns}rPr')
                fmt = _parse_run_props(rpr, wns) if rpr is not None else ""
                parts.append(f"[{style}]{fmt} {content}")

            elif tag == 'tbl':
                rows = child.findall(f'.//{wns}tr')
                parts.append(f"[TABLE: {len(rows)} rows]")
                for ri, row in enumerate(rows[:5]):
                    cells = row.findall(f'.//{wns}tc')
                    cell_texts = []
                    for cell in cells:
                        ct = ''.join(t.text or '' for t in cell.findall(f'.//{wns}t')).strip()
                        cell_texts.append(ct if ct else "___")
                    parts.append(f"  | {'  |  '.join(cell_texts)} |")
                if len(rows) > 5:
                    parts.append(f"  ... +{len(rows)-5} more rows")

    return "\n".join(parts)


def _parse_run_props(rpr, wns: str) -> str:
    """Parse run properties into a compact format string."""
    props = []
    sz = rpr.find(f'{wns}sz')
    if sz is not None:
        props.append(f"{int(sz.get(f'{wns}val'))/2}pt")
    b = rpr.find(f'{wns}b')
    if b is not None:
        props.append("bold")
    i = rpr.find(f'{wns}i')
    if i is not None:
        props.append("italic")
    fonts = rpr.find(f'{wns}rFonts')
    if fonts is not None:
        fname = fonts.get(f'{wns}ascii') or fonts.get(f'{wns}asciiTheme')
        if fname:
            props.append(f"font:{fname}")
    color = rpr.find(f'{wns}color')
    if color is not None:
        cval = color.get(f'{wns}val')
        if cval and cval != '000000':
            props.append(f"color:#{cval}")
    return f"({', '.join(props)})" if props else ""


def _bedrock_extract_structure(text: str, ext: str, model_id: str) -> dict | None:
    """Use Bedrock to analyze document text and return detailed structure."""
    import boto3

    client = boto3.client("bedrock-runtime", region_name=os.getenv("AWS_REGION", "us-east-1"))

    format_map = {".pdf": "PDF", ".docx": "Word document", ".doc": "Word document",
                  ".pptx": "PowerPoint", ".md": "Markdown"}
    doc_format = format_map.get(ext, "document")

    prompt = f"""Analyze this {doc_format} and extract its COMPLETE structure including formatting, layout, and text placement.

The text below includes metadata markers in brackets like [THEME FONTS:...], [PAGE SIZE:...], [MARGINS:...], [DEFAULT SIZE:...], [FONTS USED:...], style names in [StyleName], and formatting in (size, bold, italic, font:name).

Return a JSON object with this exact schema:
{{
  "type": "form" | "report" | "presentation" | "document",
  "title": "document title",
  "fonts": {{
    "heading_font": "font family name for headings",
    "body_font": "font family name for body text",
    "default_size_pt": 12
  }},
  "page_layout": {{
    "size": "8.5x11" or "A4" etc,
    "margins": {{"top": 0.75, "bottom": 0.44, "left": 1.5, "right": 1.5}},
    "orientation": "portrait" | "landscape"
  }},
  "sections": [
    {{
      "heading": "section heading or empty string",
      "level": 1,
      "style": "style name if known",
      "elements": [
        {{"type": "paragraph", "text": "actual text content", "style": "normal|bold|italic|centered", "font_size_pt": 12}},
        {{"type": "heading", "text": "heading text", "level": 1, "font_size_pt": 16, "bold": true}},
        {{"type": "field", "label": "field label"}},
        {{"type": "checkbox", "label": "checkbox text"}},
        {{"type": "bullet", "text": "bullet point text"}},
        {{"type": "table_header", "columns": ["col1", "col2"]}},
        {{"type": "table_row", "cells": ["cell1", "cell2"]}},
        {{"type": "signature_line"}},
        {{"type": "page_break"}},
        {{"type": "note", "text": "footnote or annotation"}}
      ]
    }}
  ],
  "source_format": "{ext.lstrip('.')}",
  "page_count": null,
  "formatting": {{
    "has_headers_footers": true/false,
    "has_page_numbers": true/false,
    "has_table_of_contents": true/false,
    "layout": "single_column" | "two_column" | "mixed"
  }}
}}

Capture section headings and structure. For tables, show only column headers and row count. For long paragraphs, truncate to first 50 chars. Keep response under 3000 tokens.

Document content:
Document content:
{text[:4000]}

{f'... [TRUNCATED MIDDLE] ...' + chr(10) + text[-2000:] if len(text) > 5000 else ''}"""

    resp = client.converse(
        modelId=model_id,
        messages=[{"role": "user", "content": [{"text": prompt}]}],
        inferenceConfig={"maxTokens": 4096},
    )

    resp_text = resp["output"]["message"]["content"][0]["text"]

    # Parse JSON from response
    try:
        start = resp_text.find("{")
        end = resp_text.rfind("}") + 1
        if start >= 0 and end > start:
            structure = json.loads(resp_text[start:end])
            # Ensure required fields
            structure.setdefault("source_format", ext.lstrip("."))
            structure.setdefault("type", "document")
            structure.setdefault("title", "Untitled")
            structure.setdefault("sections", [])
            return structure
    except json.JSONDecodeError:
        logger.warning("Failed to parse Bedrock template response as JSON")

    return None


def _fix_bibliography_pattern(structure: dict) -> None:
    """Replace '. . , .' table rows with descriptive bibliography field labels,
    and simplify index tables with fill-in blanks."""
    for s in structure.get("sections", []):
        elements = s.get("elements", [])
        table_rows = [e for e in elements if e.get("type") in ("table_row", "table_header")]
        if not table_rows:
            continue

        # Bibliography: all cells are '. . , .' or '___'
        if all(
            all(c in (". . , .", "___", "") for c in (e.get("cells") or e.get("columns") or []))
            for e in table_rows
        ):
            s["elements"] = [
                {"type": "field", "label": "[Last, First Name of Author]. [Title]. [Publisher], [year of publication]."},
                {"type": "note", "text": f"Repeats {len(table_rows)} times (fill-in bibliography entries)"},
            ]
        # Index: columns separated by empty/spacer columns with content in alternating columns
        elif any(
            any(c in ("___", "") for c in (e.get("cells") or e.get("columns") or []))
            for e in table_rows
        ):
            # Check if this is a multi-column index (alternating content/spacer pattern)
            first_content_row = next((e for e in table_rows if any(c and c not in ("___", "") for c in (e.get("cells") or e.get("columns") or []))), None)
            if first_content_row:
                first_row = first_content_row.get("cells") or first_content_row.get("columns") or []
                is_index_layout = len(first_row) >= 3 and all(
                    first_row[i] in ("", "___") for i in range(1, len(first_row), 2)
                )
            else:
                is_index_layout = False

            if is_index_layout:
                # Multi-column alphabetical index
                # Capital letters are section headers, entries below are "term, page_number"
                num_cols = (len(first_row) + 1) // 2
                new_elements = [{"type": "note", "text": f"Alphabetical index ({num_cols} columns, {len(table_rows)} rows)"}]
                for e in table_rows:
                    cells = e.get("cells") or e.get("columns") or []
                    col_entries = [cells[i] for i in range(0, len(cells), 2) if i < len(cells)]
                    row_el = {"type": "index_row", "columns": []}
                    for entry in col_entries:
                        if not entry or entry == "___":
                            row_el["columns"].append("")
                        elif len(entry) <= 2 and entry.isalpha() and entry[0].isupper():
                            row_el["columns"].append({"letter": entry})
                        else:
                            # "Term, page_number" format
                            row_el["columns"].append({"entry": entry})
                    new_elements.append(row_el)
                s["elements"] = new_elements
            else:
                # Generic table with some blanks
                sample_rows = []
                blank_count = 0
                for e in table_rows:
                    cells = e.get("cells") or e.get("columns") or []
                    non_blank = [c for c in cells if c and c != "___"]
                    if non_blank:
                        sample_rows.append(non_blank)
                    else:
                        blank_count += 1
                new_elements = []
                if sample_rows:
                    new_elements.append({"type": "note", "text": f"Table with {len(sample_rows)} entries:"})
                    for row in sample_rows[:6]:
                        new_elements.append({"type": "paragraph", "text": " | ".join(row)})
                    if len(sample_rows) > 6:
                        new_elements.append({"type": "note", "text": f"... +{len(sample_rows)-6} more"})
                if blank_count > 0:
                    new_elements.append({"type": "note", "text": f"+ {blank_count} blank fill-in rows"})
                if new_elements:
                    s["elements"] = new_elements


def _enrich_pdf_structure(structure: dict, text: str) -> None:
    """Enrich Bedrock PDF result with fill-in patterns detected from raw text."""
    # Detect bibliography fill-in pattern
    if "[Last, First Name" in text or "[Publisher]" in text:
        # Check if Bibliography section already has content
        has_bib_content = False
        for s in structure.get("sections", []):
            if "iblio" in s.get("heading", "").lower():
                if any(e.get("type") == "field" for e in s.get("elements", [])):
                    has_bib_content = True
                else:
                    # Add the fill-in description
                    s["elements"] = [
                        {"type": "field", "label": "[Last, First Name of Author]. [Title]. [Publisher], [year of publication]."},
                        {"type": "note", "text": "Repeating fill-in bibliography entries"},
                    ]
                    has_bib_content = True
                break
        if not has_bib_content:
            # Add Bibliography section if missing
            structure["sections"].append({
                "heading": "Bibliography",
                "elements": [
                    {"type": "field", "label": "[Last, First Name of Author]. [Title]. [Publisher], [year of publication]."},
                    {"type": "note", "text": "Repeating fill-in bibliography entries"},
                ],
            })
        structure["type"] = "form"
        structure.setdefault("fill_in_count", 1)

    # Fix page_layout "unknown" values from PDF metadata
    pl = structure.get("page_layout", {})
    if pl.get("size") == "unknown" or not pl.get("size"):
        # Standard US Letter assumed for most templates
        structure.setdefault("page_layout", {})
        structure["page_layout"]["size"] = "8.5x11"
        structure["page_layout"]["orientation"] = "portrait"


def _extract_pdf_template(file_bytes: bytes) -> dict:
    """Extract structure from a PDF by analyzing text layout."""
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(file_bytes))
    all_text = "\n".join((p.extract_text() or "") for p in reader.pages)

    # Parse structure from text
    sections = []
    current_section: dict = {"heading": "", "elements": []}
    title = "Untitled"

    for line in all_text.split("\n"):
        stripped = line.strip()
        if not stripped:
            continue

        # Heuristic: short uppercase or title-case lines are headings
        # Filter out addresses, names, dates, and other false positives
        import re
        _looks_like_address = bool(re.search(r'\d{2,}.*(?:Ave|St|Rd|Blvd|Dr|Ln|Way|Street|Road)', stripped, re.IGNORECASE))
        _looks_like_name = bool(re.match(r'^(?:By |Prof(?:essor)?\.? |Dr\.? )?[A-Z][a-z]+ [A-Z][a-z]+$', stripped))
        _looks_like_number = bool(re.match(r'^[\d\s/.,XxI]+$', stripped))

        is_heading = (
            not _looks_like_address
            and not _looks_like_name
            and not _looks_like_number
            and (
                (stripped.isupper() and 3 < len(stripped) < 80 and len(stripped.split()) >= 2 and len(stripped.split()) <= 8)
                or (stripped.istitle() and len(stripped) > 10 and len(stripped) < 60 and len(stripped.split()) >= 3
                    and not stripped.startswith("By "))
            )
        )

        if is_heading and len(stripped) > 2:
            if current_section["heading"] or current_section["elements"]:
                sections.append(current_section)
            current_section = {"heading": stripped, "elements": []}
            if title == "Untitled":
                title = stripped
        elif stripped.startswith("☐") or stripped.startswith("□"):
            current_section["elements"].append({"type": "checkbox", "label": stripped[1:].strip()})
        elif "___" in stripped or "____" in stripped:
            label = stripped.split("___")[0].strip().rstrip(":")
            current_section["elements"].append({"type": "field", "label": label or "Blank"})
        elif stripped.startswith("•") or stripped.startswith("-") or stripped.startswith("●"):
            current_section["elements"].append({"type": "bullet", "text": stripped.lstrip("•-● ").strip()[:100]})
        elif len(stripped) < 120:
            current_section["elements"].append({"type": "paragraph", "text": stripped[:120]})

        # Cap sections to keep preview manageable
        if len(sections) >= 20:
            break

    if current_section["heading"] or current_section["elements"]:
        sections.append(current_section)

    has_fields = any(
        e.get("type") in ("field", "checkbox") for s in sections for e in s.get("elements", [])
    )

    return {
        "type": "form" if has_fields else "document",
        "title": title,
        "sections": sections[:20],
        "source_format": "pdf",
        "page_count": len(reader.pages),
    }


def _extract_docx_template(file_bytes: bytes) -> dict:
    """Extract structure from a DOCX by walking the XML tree directly."""
    import zipfile
    from lxml import etree

    ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
          'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    wns = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'

    fonts = {}
    page_layout = {}
    sections = []
    current_section: dict = {"heading": "", "elements": []}
    fill_in_count = 0

    with zipfile.ZipFile(BytesIO(file_bytes)) as z:
        # Theme fonts
        for tf in [n for n in z.namelist() if 'theme' in n and n.endswith('.xml')]:
            try:
                theme_root = etree.fromstring(z.read(tf))
                major = theme_root.find('.//a:majorFont/a:latin', ns)
                minor = theme_root.find('.//a:minorFont/a:latin', ns)
                if major is not None:
                    fonts["heading_font"] = major.get('typeface')
                if minor is not None:
                    fonts["body_font"] = minor.get('typeface')
            except Exception:
                pass

        # Default size from styles
        for sf in [n for n in z.namelist() if 'styles' in n and n.endswith('.xml') and 'glossary' not in n]:
            try:
                styles_root = etree.fromstring(z.read(sf))
                sz = styles_root.find(f'.//{wns}docDefaults//{wns}sz')
                if sz is not None:
                    fonts["default_size_pt"] = int(sz.get(f'{wns}val')) / 2
            except Exception:
                pass

        # Document body
        doc_xml = z.read('word/document.xml')
        root = etree.fromstring(doc_xml)

        # Page layout
        sect_pr = root.find(f'.//{wns}sectPr')
        if sect_pr is not None:
            pg_sz = sect_pr.find(f'{wns}pgSz')
            pg_mar = sect_pr.find(f'{wns}pgMar')
            if pg_sz is not None:
                w = int(pg_sz.get(f'{wns}w', 0)) / 1440
                h = int(pg_sz.get(f'{wns}h', 0)) / 1440
                page_layout["size"] = f"{w:.1f}x{h:.1f}"
                page_layout["orientation"] = "landscape" if w > h else "portrait"
            if pg_mar is not None:
                page_layout["margins"] = {
                    k: round(int(pg_mar.get(f'{wns}{k}', 0)) / 1440, 2)
                    for k in ['top', 'bottom', 'left', 'right']
                }

        # Walk body
        body = root.find(f'{wns}body')
        if body is None:
            return {"type": "document", "title": "Untitled", "sections": [], "source_format": "docx"}

        for child in body:
            tag = child.tag.split('}')[-1]

            if tag == 'sdt':
                texts = child.findall(f'.//{wns}t')
                content = ''.join(t.text or '' for t in texts).strip()
                if not content:
                    continue

                # Check if placeholder (fill-in field)
                sdt_pr = child.find(f'{wns}sdtPr')
                is_placeholder = sdt_pr is not None and sdt_pr.find(f'{wns}showingPlcHdr') is not None

                # Get style
                pstyle = child.find(f'.//{wns}pStyle')
                style = pstyle.get(f'{wns}val') if pstyle is not None else ""

                # Get formatting
                rpr = child.find(f'.//{wns}rPr')
                fmt = {}
                if rpr is not None:
                    sz_el = rpr.find(f'{wns}sz')
                    if sz_el is not None:
                        fmt["font_size_pt"] = int(sz_el.get(f'{wns}val')) / 2
                    if rpr.find(f'{wns}b') is not None:
                        fmt["bold"] = True
                    if rpr.find(f'{wns}i') is not None:
                        fmt["italic"] = True

                # Detect headings by style
                if 'Heading' in style or 'UnnumberedHeading' in style or 'TOCHeading' in style:
                    if current_section["heading"] or current_section["elements"]:
                        sections.append(current_section)
                    current_section = {"heading": content, "style": style, "elements": []}
                elif is_placeholder:
                    fill_in_count += 1
                    current_section["elements"].append({"type": "field", "label": content[:80], "style": style, **fmt})
                else:
                    current_section["elements"].append({"type": "paragraph", "text": content[:120], "style": style, **fmt})

            elif tag == 'p':
                texts = child.findall(f'.//{wns}t')
                content = ''.join(t.text or '' for t in texts).strip()
                if not content:
                    continue
                pstyle = child.find(f'.//{wns}pStyle')
                style = pstyle.get(f'{wns}val') if pstyle is not None else ""
                if 'Heading' in style:
                    if current_section["heading"] or current_section["elements"]:
                        sections.append(current_section)
                    current_section = {"heading": content, "style": style, "elements": []}
                elif '___' in content:
                    fill_in_count += 1
                    label = content.split("___")[0].strip().rstrip(":")
                    current_section["elements"].append({"type": "field", "label": label or "Blank"})
                else:
                    current_section["elements"].append({"type": "paragraph", "text": content[:120], "style": style})

            elif tag == 'tbl':
                rows = child.findall(f'.//{wns}tr')
                # If current section has a heading but no elements, attach table to it
                if current_section["heading"] and not current_section["elements"]:
                    table_section = current_section
                else:
                    if current_section["heading"] or current_section["elements"]:
                        sections.append(current_section)
                    table_section = {"heading": f"Table ({len(rows)} rows)", "elements": []}
                has_blanks = False
                for ri, row in enumerate(rows):
                    cells = row.findall(f'.//{wns}tc')
                    cell_texts = []
                    for cell in cells:
                        ct = ''.join(t.text or '' for t in cell.findall(f'.//{wns}t')).strip()
                        cell_texts.append(ct if ct else "___")
                        if not ct:
                            has_blanks = True
                    if any(c != "___" for c in cell_texts):
                        table_section["elements"].append({"type": "table_row", "cells": cell_texts})
                if has_blanks:
                    fill_in_count += len(rows)
                    table_section["elements"].insert(0, {"type": "note", "text": f"Fill-in table with {len(rows)} rows"})
                # Detect bibliography pattern: ". . , ." = [Author]. [Title]. [Publisher], [year].
                if all(c in (". . , .", "___") for row_el in table_section.get("elements", []) if row_el.get("type") == "table_row" for c in row_el.get("cells", [])):
                    table_section["elements"] = [
                        {"type": "field", "label": "[Last, First Name of Author]. [Title]. [Publisher], [year of publication]."},
                        {"type": "note", "text": f"Repeats {len(rows)} times (fill-in bibliography entries)"},
                    ]
                table_section["row_count"] = len(rows)
                sections.append(table_section)
                current_section = {"heading": "", "elements": []}

    if current_section["heading"] or current_section["elements"]:
        sections.append(current_section)

    # Determine type
    doc_type = "form" if fill_in_count > 0 else "document"

    # Derive title from first meaningful content
    title = "Untitled"
    # First check table rows for a title-like entry (first non-blank cell in first table)
    for s in sections:
        if s.get("row_count") and s.get("elements"):
            for e in s["elements"]:
                if e.get("type") == "table_row":
                    for cell in (e.get("cells") or []):
                        if cell and cell != "___" and len(cell) > 2:
                            title = cell
                            break
                if title != "Untitled":
                    break
        if title != "Untitled":
            break
    # Fallback to headings/paragraphs
    if title == "Untitled":
        for s in sections:
            if s.get("heading") and not s["heading"].startswith("Table "):
                title = s["heading"]
                break

    return {
        "type": doc_type,
        "title": title,
        "sections": sections,
        "source_format": "docx",
        "fonts": fonts if fonts else None,
        "page_layout": page_layout if page_layout else None,
        "table_count": len([c for c in body if c.tag.split('}')[-1] == 'tbl']),
        "fill_in_count": fill_in_count,
        "formatting": {
            "has_table_of_contents": any('TOC' in s.get('style', '') for s in sections),
            "has_headers_footers": sect_pr is not None and sect_pr.find(f'{wns}headerReference') is not None,
        },
    }


def _extract_pptx_template(file_bytes: bytes) -> dict:
    """Extract structure from a PPTX."""
    from pptx import Presentation

    prs = Presentation(BytesIO(file_bytes))
    slides = []

    for slide in prs.slides:
        slide_data = {"title": "", "bullets": [], "has_image": False}

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if not slide_data["title"]:
                        slide_data["title"] = text
                    else:
                        slide_data["bullets"].append(text)
            if shape.shape_type == 13:  # Picture
                slide_data["has_image"] = True

        slides.append(slide_data)

    return {
        "type": "presentation",
        "title": slides[0]["title"] if slides else "Untitled",
        "slide_count": len(slides),
        "slides": [
            {
                "title": s["title"],
                "bullet_count": len(s["bullets"]),
                "has_image": s["has_image"],
            }
            for s in slides
        ],
        "source_format": "pptx",
    }


def _extract_md_template(text: str) -> dict:
    """Extract structure from Markdown text."""
    sections = []
    current_section = {"heading": "", "level": 0, "elements": []}

    for line in text.split("\n"):
        stripped = line.strip()

        if stripped.startswith("#"):
            if current_section["heading"] or current_section["elements"]:
                sections.append(current_section)
            level = len(stripped) - len(stripped.lstrip("#"))
            heading = stripped.lstrip("#").strip()
            current_section = {"heading": heading, "level": level, "elements": []}
        elif stripped.startswith("- ") or stripped.startswith("* "):
            current_section["elements"].append({"type": "bullet", "text": stripped[2:][:80]})
        elif stripped.startswith("|"):
            current_section["elements"].append({"type": "table_row"})
        elif "___" in stripped:
            label = stripped.split("___")[0].strip().rstrip(":")
            current_section["elements"].append({"type": "field", "label": label})
        elif stripped:
            current_section["elements"].append({"type": "paragraph", "text": stripped[:80]})

    if current_section["heading"] or current_section["elements"]:
        sections.append(current_section)

    has_fields = any(
        e["type"] == "field" for s in sections for e in s.get("elements", [])
    )

    return {
        "type": "form" if has_fields else "document",
        "title": sections[0]["heading"] if sections else "Untitled",
        "sections": sections,
        "source_format": "md",
    }


def _extract_image_template(file_bytes: bytes, ext: str) -> dict:
    """Send image to vision LLM to describe the document structure."""
    import boto3

    model_id = os.getenv("BEDROCK_VISION_MODEL_ID", "")
    if not model_id:
        raise ValueError("No Vision OCR model configured")

    fmt_map = {".jpg": "jpeg", ".jpeg": "jpeg", ".png": "png", ".tiff": "tiff", ".tif": "tiff"}
    fmt = fmt_map.get(ext, "jpeg")

    client = boto3.client("bedrock-runtime", region_name=os.getenv("AWS_REGION", "us-east-1"))

    resp = client.converse(
        modelId=model_id,
        messages=[{
            "role": "user",
            "content": [
                {"image": {"format": fmt, "source": {"bytes": file_bytes}}},
                {"text": (
                    "Analyze this document/form image and describe its STRUCTURE only (not content). "
                    "Return a JSON object with: "
                    '{"type": "form"|"document"|"report", "title": "...", "sections": [{"heading": "...", "elements": [{"type": "field"|"checkbox"|"paragraph"|"signature_line"|"table", "label": "..."}]}]}'
                    " Be precise about field labels, section headings, and element types."
                )},
            ],
        }],
        inferenceConfig={"maxTokens": 2048},
    )

    text = resp["output"]["message"]["content"][0]["text"]

    # Try to parse as JSON
    try:
        # Find JSON in the response
        start = text.find("{")
        end = text.rfind("}") + 1
        if start >= 0 and end > start:
            structure = json.loads(text[start:end])
            structure["source_format"] = "image"
            return structure
    except json.JSONDecodeError:
        pass

    # Fallback: return raw description
    return {
        "type": "document",
        "title": "Extracted from image",
        "description": text,
        "sections": [],
        "source_format": "image",
    }


def _ask_bedrock_for_structure(text: str, source_type: str) -> dict:
    """Ask Bedrock to analyze document text and return structure JSON."""
    import boto3

    model_id = os.getenv("BEDROCK_MODEL_ID", "")
    if not model_id:
        raise ValueError("No model configured")

    client = boto3.client("bedrock-runtime", region_name=os.getenv("AWS_REGION", "us-east-1"))

    resp = client.converse(
        modelId=model_id,
        system=[{"text": (
            "You analyze document structure. Return ONLY a JSON object describing the layout. "
            "Do not include the actual content, just the structure pattern."
        )}],
        messages=[{
            "role": "user",
            "content": [{"text": (
                f"Analyze this {source_type} document and describe its structure as JSON:\n\n"
                f"{text[:3000]}\n\n"
                "Return: {\"type\": \"form\"|\"document\"|\"report\", \"title\": \"...\", "
                "\"sections\": [{\"heading\": \"...\", \"elements\": [{\"type\": \"field\"|\"checkbox\"|\"paragraph\"|\"signature_line\"|\"table\", \"label\": \"...\"}]}]}"
            )}],
        }],
        inferenceConfig={"maxTokens": 2048},
    )

    text_resp = resp["output"]["message"]["content"][0]["text"]

    try:
        start = text_resp.find("{")
        end = text_resp.rfind("}") + 1
        if start >= 0 and end > start:
            return json.loads(text_resp[start:end])
    except json.JSONDecodeError:
        pass

    return {"type": "document", "title": "Untitled", "sections": []}
