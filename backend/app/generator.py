"""Document generation from user prompts + indexed document context.

Uses RAG (retrieval-augmented generation) to create documents grounded
in the user's actual house documents. Supports multiple output formats:
  - Markdown (.md) - direct from Bedrock
  - Word (.docx) - via Pandoc with optional reference template
  - PDF (.pdf) - via Pandoc + weasyprint
  - Image (.png) - Markdown to styled HTML, screenshot via Playwright
  - PowerPoint (.pptx) - via Pandoc slide-per-heading
"""

from __future__ import annotations

import logging
import os
import subprocess
import tempfile
from pathlib import Path

logger = logging.getLogger(__name__)

_BASE_RULES = (
    "- Use only information from the provided source documents\n"
    "- Include specific details, numbers, addresses, and names from the source material\n"
    "- If the source material doesn't contain enough information, say so clearly\n"
    "- Write in plain English that a homeowner would understand"
)

_FORMAT_INSTRUCTIONS = {
    "md": (
        "Structure the output as clean Markdown with proper headings (# ## ###), "
        "bullet points, and paragraphs. Use bold for key terms."
    ),
    "docx": (
        "Structure the output as a formal document with:\n"
        "- A clear title as a level-1 heading\n"
        "- Sections with level-2 headings\n"
        "- Proper paragraphs (not just bullet points)\n"
        "- Tables where data comparison is useful\n"
        "- A professional, complete tone suitable for a Word document"
    ),
    "pdf": (
        "Structure the output as a formal report with:\n"
        "- A title and subtitle\n"
        "- Sections with clear headings\n"
        "- Detailed paragraphs with supporting data\n"
        "- Tables for structured information\n"
        "- A conclusion or summary section"
    ),
    "png": (
        "Structure the output as a concise, single-page reference card:\n"
        "- Keep it brief and scannable\n"
        "- Use short bullet points, not long paragraphs\n"
        "- Organize into 2-4 clear sections\n"
        "- Include only the most important facts\n"
        "- Aim for no more than 30 lines total"
    ),
    "pptx": (
        "Structure the output as a presentation with EXACTLY this format:\n"
        "- First level-1 heading (# Title) becomes the title slide\n"
        "- Each level-2 heading (## Slide Title) becomes a new slide\n"
        "- Under each slide heading, use bullet points (- item)\n"
        "- Keep each slide to 4-6 bullet points maximum\n"
        "- Keep bullet points short (one line each)\n"
        "- Include a final slide with summary or contact info\n"
        "- Aim for 5-8 slides total"
    ),
    "form": (
        "Structure the output as a fillable form with:\n"
        "- A form title as level-1 heading\n"
        "- Sections with level-2 headings\n"
        "- Field labels followed by blanks: **Field Name:** _______________\n"
        "- Checkbox items: ☐ Option text\n"
        "- Signature lines: **Signature:** ___________________________ **Date:** ___________\n"
        "- Pre-fill any fields where the information is available in the source documents\n"
        "- Include all addresses, phone numbers, and submission instructions from the source"
    ),
}


def generate_markdown(prompt: str, context: str, manual_mode: bool = False, fmt: str = "md") -> str:
    """Ask Bedrock to generate a markdown document from context and prompt."""
    import boto3

    model_id = os.getenv("BEDROCK_GENERATE_MODEL_ID", "") or os.getenv("BEDROCK_MODEL_ID", "")
    if not model_id:
        raise ValueError("No Ask AI model configured. Set it in Settings.")

    client = boto3.client(
        "bedrock-runtime",
        region_name=os.getenv("AWS_REGION", "us-east-1"),
    )

    # Detect if this is a form request
    effective_fmt = fmt
    form_keywords = ["form", "application", "fill out", "fill in", "request form"]
    if any(kw in prompt.lower() for kw in form_keywords):
        effective_fmt = "form"

    format_instruction = _FORMAT_INSTRUCTIONS.get(effective_fmt, _FORMAT_INSTRUCTIONS["md"])

    mode_intro = (
        "The user has selected specific documents as source material. "
        "Their prompt describes WHAT to create, not what to search for."
        if manual_mode
        else "Your job is to create well-structured documents using ONLY the provided source material."
    )

    system_prompt = (
        f"You are a professional document writer. {mode_intro}\n\n"
        f"Output format instructions:\n{format_instruction}\n\n"
        f"Rules:\n{_BASE_RULES}"
    )

    user_msg = (
        f"Source documents:\n{context}\n\n"
        f"Request: {prompt}\n\n"
        "Generate the document in Markdown format following the format instructions above."
    )

    resp = client.converse(
        modelId=model_id,
        system=[{"text": system_prompt}],
        messages=[{"role": "user", "content": [{"text": user_msg}]}],
        inferenceConfig={"maxTokens": 4096},
    )

    content = resp["output"]["message"]["content"][0]["text"]

    # Track usage
    if os.getenv("TRACK_USAGE", "true").lower() == "true":
        usage = resp.get("usage", {})
        try:
            from .pricing import estimate_cost
            from .db import get_conn

            cost = estimate_cost(
                model_id,
                usage.get("inputTokens", 0),
                usage.get("outputTokens", 0),
                os.getenv("AWS_REGION", "us-east-1"),
            )
            conn = get_conn()
            with conn.cursor() as cur:
                cur.execute(
                    """INSERT INTO token_usage
                       (model_id, operation, input_tokens, output_tokens, estimated_cost_usd)
                       VALUES (%s, %s, %s, %s, %s)""",
                    (model_id, "generate", usage.get("inputTokens", 0),
                     usage.get("outputTokens", 0), cost),
                )
            conn.close()
        except Exception:
            pass

    return content


def convert_to_docx(markdown_content: str) -> bytes:
    """Convert markdown to a styled DOCX using python-docx for full control."""
    from docx import Document
    from docx.shared import Pt, Inches, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.style import WD_STYLE_TYPE
    import io

    doc = Document()

    # Set margins
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    # Style the default font
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
    style.paragraph_format.space_after = Pt(6)
    style.paragraph_format.line_spacing = 1.15

    # Style headings
    for level in range(1, 4):
        h_style = doc.styles[f"Heading {level}"]
        h_style.font.name = "Calibri"
        h_style.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
        if level == 1:
            h_style.font.size = Pt(24)
            h_style.font.bold = True
            h_style.paragraph_format.space_before = Pt(0)
            h_style.paragraph_format.space_after = Pt(12)
        elif level == 2:
            h_style.font.size = Pt(16)
            h_style.font.bold = True
            h_style.paragraph_format.space_before = Pt(18)
            h_style.paragraph_format.space_after = Pt(6)
        else:
            h_style.font.size = Pt(13)
            h_style.font.bold = True

    # Parse and build document
    lines = markdown_content.split("\n")
    for line in lines:
        stripped = line.rstrip()

        if stripped.startswith("# "):
            doc.add_heading(stripped[2:].strip(), level=1)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:].strip(), level=2)
        elif stripped.startswith("### "):
            doc.add_heading(stripped[4:].strip(), level=3)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            p = doc.add_paragraph(stripped[2:].strip(), style="List Bullet")
            p.paragraph_format.space_after = Pt(4)
        elif stripped.startswith("☐ "):
            p = doc.add_paragraph("☐ " + stripped[2:].strip())
            p.paragraph_format.left_indent = Cm(1)
            p.paragraph_format.space_after = Pt(4)
        elif stripped.startswith("**") and stripped.endswith("**"):
            p = doc.add_paragraph()
            run = p.add_run(stripped.strip("*").strip())
            run.bold = True
        elif "___" in stripped or "—" * 5 in stripped:
            # Signature/form line
            p = doc.add_paragraph(stripped)
            p.paragraph_format.space_before = Pt(12)
        elif stripped.strip():
            # Handle inline bold
            p = doc.add_paragraph()
            _add_rich_text(p, stripped)
        # Skip empty lines (they become paragraph spacing)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _add_rich_text(paragraph, text: str):
    """Parse simple markdown bold (**text**) into Word runs."""
    import re
    parts = re.split(r"(\*\*[^*]+\*\*)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            run = paragraph.add_run(part[2:-2])
            run.bold = True
        else:
            paragraph.add_run(part)


def convert_to_pdf(markdown_content: str) -> bytes:
    """Convert markdown to PDF via DOCX intermediate for consistent styling."""
    # Generate styled DOCX first
    docx_bytes = convert_to_docx(markdown_content)

    # Convert DOCX to PDF using weasyprint via HTML
    # (since LibreOffice may not be available in container)
    import markdown as md_lib
    from weasyprint import HTML

    html = md_lib.markdown(markdown_content, extensions=["tables", "fenced_code"])
    styled = f"""<html><head><style>
        @page {{ margin: 2.5cm; size: A4; }}
        body {{ font-family: Calibri, sans-serif; font-size: 11pt; color: #374151; line-height: 1.5; }}
        h1 {{ font-size: 24pt; color: #1a1a2e; font-weight: bold; margin-top: 0; margin-bottom: 12pt; }}
        h2 {{ font-size: 16pt; color: #1a1a2e; font-weight: bold; margin-top: 18pt; margin-bottom: 6pt;
              border-bottom: 2px solid #6366f1; padding-bottom: 4pt; }}
        h3 {{ font-size: 13pt; color: #1a1a2e; font-weight: bold; }}
        ul {{ padding-left: 20pt; }}
        li {{ margin-bottom: 4pt; }}
        table {{ border-collapse: collapse; width: 100%; margin: 12pt 0; }}
        td, th {{ border: 1px solid #e5e7eb; padding: 8pt 12pt; font-size: 10pt; }}
        th {{ background: #f3f4f6; font-weight: bold; }}
        p {{ margin-bottom: 6pt; }}
    </style></head><body>{html}</body></html>"""

    return HTML(string=styled).write_pdf()


def convert_to_png(markdown_content: str) -> bytes:
    """Convert markdown to a styled PNG image using Playwright."""
    import markdown as md_lib
    from playwright.sync_api import sync_playwright

    html = md_lib.markdown(markdown_content, extensions=["tables", "fenced_code"])
    styled = f"""<html><head><style>
        body {{ font-family: -apple-system, sans-serif; max-width: 800px; margin: 0 auto;
               padding: 40px; line-height: 1.6; color: #1a1a2e; background: #fff; }}
        h1 {{ font-size: 1.8rem; color: #1a1a2e; border-bottom: 2px solid #6366f1; padding-bottom: 8px; }}
        h2 {{ font-size: 1.3rem; color: #374151; margin-top: 24px; }}
        ul, ol {{ padding-left: 24px; }}
        table {{ border-collapse: collapse; width: 100%; margin: 16px 0; }}
        td, th {{ border: 1px solid #e5e7eb; padding: 8px 12px; font-size: 0.9rem; }}
        th {{ background: #f3f4f6; font-weight: 600; }}
        blockquote {{ border-left: 3px solid #6366f1; padding-left: 12px; color: #6b7280; }}
        code {{ background: #f3f4f6; padding: 2px 6px; border-radius: 4px; font-size: 0.85rem; }}
    </style></head><body>{html}</body></html>"""

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 900, "height": 600})
        page.set_content(styled)
        page.wait_for_load_state("networkidle")
        png_bytes = page.screenshot(full_page=True)
        browser.close()

    return png_bytes


def convert_to_pptx(markdown_content: str) -> bytes:
    """Convert markdown to a styled PPTX using python-pptx for full control."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    import io

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Parse markdown into slides
    slides_data = _parse_slides(markdown_content)

    # Colors
    NAVY = RGBColor(0x1A, 0x1A, 0x2E)
    INDIGO = RGBColor(0x63, 0x66, 0xF1)
    GRAY = RGBColor(0x4B, 0x55, 0x63)
    LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    for i, slide_data in enumerate(slides_data):
        if i == 0:
            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            # Dark background
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = NAVY

            # Title
            left, top = Inches(1), Inches(2.5)
            txBox = slide.shapes.add_textbox(left, top, Inches(11), Inches(2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data["title"]
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

            # Subtitle line
            if slide_data.get("bullets"):
                p2 = tf.add_paragraph()
                p2.text = slide_data["bullets"][0] if slide_data["bullets"] else ""
                p2.font.size = Pt(18)
                p2.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
                p2.alignment = PP_ALIGN.CENTER
        else:
            # Content slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = LIGHT_BG

            # Title bar
            title_bar = slide.shapes.add_shape(
                1, Inches(0), Inches(0), prs.slide_width, Inches(1.2),
            )
            title_bar.fill.solid()
            title_bar.fill.fore_color.rgb = NAVY
            title_bar.line.fill.background()

            # Title text
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11), Inches(0.9))
            tf = txBox.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = slide_data["title"]
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = WHITE

            # Bullet points
            if slide_data.get("bullets"):
                content_box = slide.shapes.add_textbox(
                    Inches(1), Inches(1.6), Inches(11), Inches(5.2),
                )
                tf = content_box.text_frame
                tf.word_wrap = True

                for j, bullet in enumerate(slide_data["bullets"]):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = bullet
                    p.font.size = Pt(18)
                    p.font.color.rgb = GRAY
                    p.space_after = Pt(12)
                    p.level = 0

            # Slide number
            num_box = slide.shapes.add_textbox(
                Inches(12.2), Inches(7), Inches(0.8), Inches(0.4),
            )
            num_tf = num_box.text_frame
            num_p = num_tf.paragraphs[0]
            num_p.text = str(i + 1)
            num_p.font.size = Pt(10)
            num_p.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
            num_p.alignment = PP_ALIGN.RIGHT

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _parse_slides(markdown: str) -> list[dict]:
    """Parse markdown into slide data: [{title, bullets}, ...]"""
    slides = []
    current: dict = {"title": "", "bullets": []}

    for line in markdown.split("\n"):
        line = line.rstrip()
        if line.startswith("# "):
            if current["title"]:
                slides.append(current)
            current = {"title": line[2:].strip(), "bullets": []}
        elif line.startswith("## "):
            if current["title"]:
                slides.append(current)
            current = {"title": line[3:].strip(), "bullets": []}
        elif line.startswith("- ") or line.startswith("* "):
            current["bullets"].append(line[2:].strip())
        elif line.startswith("☐ "):
            current["bullets"].append("☐ " + line[2:].strip())
        elif line.startswith("**") and line.endswith("**"):
            current["bullets"].append(line.strip("*").strip())
        elif line.strip() and not line.startswith("#"):
            # Regular text becomes a bullet
            if line.strip():
                current["bullets"].append(line.strip())

    if current["title"]:
        slides.append(current)

    return slides if slides else [{"title": "Untitled", "bullets": ["No content generated"]}]
