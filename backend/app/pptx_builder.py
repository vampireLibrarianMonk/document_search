"""Enhanced PPTX builder with PlantUML diagram embedding and speaker notes.

Builds presentation decks programmatically using python-pptx with support for:
  - Embedded PlantUML diagrams (rendered to PNG, inserted via add_picture)
  - Speaker notes per slide
  - Navy/white professional theme
  - Title slides, content slides, and diagram slides
"""

from __future__ import annotations

import io
import logging
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .plantuml import is_available as plantuml_available
from .plantuml import render_puml_to_png

logger = logging.getLogger(__name__)

# Theme colors
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x4B, 0x55, 0x63)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x63, 0x66, 0xF1)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def build_pptx(slides_data: list[dict]) -> bytes:
    """Build a PPTX from structured slide data.

    Each slide dict can have:
        title: str
        bullets: list[str] (optional)
        notes: str (optional) - speaker notes
        diagram: str (optional) - PlantUML source to render and embed
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for i, sd in enumerate(slides_data):
        if i == 0:
            _add_title_slide(prs, sd)
        elif sd.get("diagram"):
            _add_diagram_slide(prs, sd, i)
        else:
            _add_content_slide(prs, sd, i)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_title_slide(prs: Presentation, sd: dict):
    """Add a dark-background title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = sd.get("title", "")
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle from first bullet
    if sd.get("bullets"):
        p2 = tf.add_paragraph()
        subtitle = sd["bullets"][0]
        p2.text = subtitle[:80] + "..." if len(subtitle) > 80 else subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        p2.alignment = PP_ALIGN.CENTER

    # Speaker notes
    if sd.get("notes"):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = sd["notes"]


def _add_content_slide(prs: Presentation, sd: dict, index: int):
    """Add a content slide with title bar and bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT_BG

    # Title bar
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY
    title_bar.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11), Inches(0.9))
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = sd.get("title", "")
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Bullets
    if sd.get("bullets"):
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11), Inches(5.2))
        tf = content_box.text_frame
        tf.word_wrap = True

        # Auto-scale: reduce font size if too many bullets or long text
        bullets = sd["bullets"]
        total_chars = sum(len(b) for b in bullets)
        box_h_inches = 5.2
        if len(bullets) > 8 or total_chars > 600:
            font_size = Pt(14)
            spacing = Pt(6)
        elif len(bullets) > 6 or total_chars > 400:
            font_size = Pt(16)
            spacing = Pt(8)
        else:
            font_size = Pt(18)
            spacing = Pt(12)

        for j, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            # Truncate extremely long bullets to prevent overflow
            p.text = bullet[:120] + "..." if len(bullet) > 120 else bullet
            p.font.size = font_size
            p.font.color.rgb = GRAY
            p.space_after = spacing

    # Slide number
    _add_slide_number(slide, index)

    # Speaker notes
    if sd.get("notes"):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = sd["notes"]


def _add_diagram_slide(prs: Presentation, sd: dict, index: int):
    """Add a slide with an embedded PlantUML diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT_BG

    # Title bar
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY
    title_bar.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11), Inches(0.9))
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = sd.get("title", "")
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Render and embed diagram
    puml_src = sd["diagram"]
    png_bytes = render_puml_to_png(puml_src)

    if png_bytes:
        # Calculate image dimensions to fit available area
        available_w = Inches(11.5)
        available_h = Inches(5.5)
        img = Image.open(io.BytesIO(png_bytes))
        img_w, img_h = img.size

        # Scale to fit
        scale = min(
            Emu(available_w) / Emu(int(img_w * 914400 / 96)),  # px to EMU at 96dpi
            Emu(available_h) / Emu(int(img_h * 914400 / 96)),
        )
        final_w = int(img_w * 914400 / 96 * scale)
        final_h = int(img_h * 914400 / 96 * scale)

        # Center the image
        left = (Emu(SLIDE_WIDTH) - final_w) // 2
        top = Inches(1.4) + (Emu(available_h) - final_h) // 2

        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp.write(png_bytes)
            tmp_path = tmp.name

        slide.shapes.add_picture(tmp_path, left, top, Emu(final_w), Emu(final_h))
        Path(tmp_path).unlink(missing_ok=True)
    else:
        # Fallback: show error text
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "[Diagram rendering failed — PlantUML not available]"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)

        # Show source as fallback
        if sd.get("bullets"):
            for bullet in sd["bullets"]:
                p2 = tf.add_paragraph()
                p2.text = bullet
                p2.font.size = Pt(14)
                p2.font.color.rgb = GRAY

    # Slide number
    _add_slide_number(slide, index)

    # Speaker notes
    if sd.get("notes"):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = sd["notes"]


def _add_slide_number(slide, index: int):
    """Add slide number in bottom-right corner."""
    num_box = slide.shapes.add_textbox(Inches(12.2), Inches(7), Inches(0.8), Inches(0.4))
    num_tf = num_box.text_frame
    num_p = num_tf.paragraphs[0]
    num_p.text = str(index + 1)
    num_p.font.size = Pt(10)
    num_p.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
    num_p.alignment = PP_ALIGN.RIGHT


def parse_enhanced_markdown(markdown: str) -> list[dict]:
    """Parse enhanced markdown with diagram blocks and notes into slide data.

    Supports:
        # Title → title slide
        ## Slide Title → content slide
        - bullet → bullet point
        ```plantuml ... ``` → diagram slide
        <!-- notes: ... --> → speaker notes for preceding slide
    """
    slides: list[dict] = []
    current: dict = {"title": "", "bullets": [], "notes": "", "diagram": ""}
    in_diagram = False
    diagram_lines: list[str] = []

    for line in markdown.split("\n"):
        stripped = line.rstrip()

        # Diagram block start
        if stripped.startswith("```plantuml") or stripped.startswith("```puml"):
            in_diagram = True
            diagram_lines = []
            continue

        # Diagram block end
        if in_diagram and stripped == "```":
            in_diagram = False
            current["diagram"] = "\n".join(diagram_lines)
            continue

        # Inside diagram block
        if in_diagram:
            diagram_lines.append(line)
            continue

        # Speaker notes (HTML comment)
        if stripped.startswith("<!-- notes:") and stripped.endswith("-->"):
            current["notes"] = stripped[11:-3].strip()
            continue
        if stripped.startswith("<!-- notes:"):
            current["notes"] = stripped[11:].strip()
            continue

        # Headings
        if stripped.startswith("# ") and not stripped.startswith("## "):
            if current["title"]:
                slides.append(current)
            current = {"title": stripped[2:].strip(), "bullets": [], "notes": "", "diagram": ""}
        elif stripped.startswith("## "):
            if current["title"]:
                slides.append(current)
            current = {"title": stripped[3:].strip(), "bullets": [], "notes": "", "diagram": ""}
        elif stripped.startswith("- ") or stripped.startswith("* "):
            current["bullets"].append(stripped[2:].strip())
        elif stripped.strip() and not stripped.startswith("#"):
            # Regular text becomes a bullet
            current["bullets"].append(stripped.strip())

    if current["title"]:
        slides.append(current)

    return slides if slides else [{"title": "Untitled", "bullets": ["No content"], "notes": "", "diagram": ""}]
